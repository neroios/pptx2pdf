#!/usr/bin/env python3
"""Convert between PPTX and PDF using LibreOffice."""

import sys
import os
import shutil
import subprocess
import glob
import tempfile
from pathlib import Path


TARGET_EXT = {".pptx": "pdf", ".pdf": "pptx"}
FAITHFUL_DPI = 300


def find_libreoffice() -> str:
    candidates = ["libreoffice", "soffice"]
    if sys.platform == "win32":
        candidates = [
            "soffice.exe",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
    for c in candidates:
        found = shutil.which(c)
        if found:
            return found
    print("Error: LibreOffice not found. Install it from https://libreoffice.org", file=sys.stderr)
    sys.exit(1)


def get_default_dir() -> str:
    if sys.platform == "win32":
        return str(Path.home() / "Documents")
    docs = os.path.expanduser("~/Documents")
    if os.path.isdir(docs):
        return docs
    return os.path.dirname(os.path.abspath(".")) or "."


def resolve_output(input_path: str, output_name: str | None, output_dir: str | None) -> str:
    in_ext = os.path.splitext(input_path)[1].lower()
    out_ext = TARGET_EXT.get(in_ext, "pdf")

    out_dir = output_dir if output_dir else get_default_dir()
    out_name = output_name if output_name else os.path.splitext(os.path.basename(input_path))[0] + "." + out_ext

    if os.path.isabs(out_name):
        return out_name

    return os.path.join(out_dir, out_name)


def convert_single(input_path: str, output_path: str, fresh: bool = False):
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    output_dir = os.path.dirname(output_path)
    os.makedirs(output_dir, exist_ok=True)

    in_ext = os.path.splitext(input_path)[1].lower()
    out_ext = os.path.splitext(output_path)[1].lower().lstrip(".")

    if out_ext not in ("pdf", "pptx"):
        print(f"Error: unsupported output format: {out_ext}", file=sys.stderr)
        sys.exit(1)

    lo = find_libreoffice()

    convert_to = out_ext
    infilter = None
    if in_ext == ".pdf" and out_ext == "pptx":
        convert_to = "pptx:Impress Office Open XML"
        infilter = "impress_pdf_import"

    cmd = [lo, "--headless", "--convert-to", convert_to, "--outdir", output_dir]
    if fresh:
        cmd += ["-env:UserInstallation=file://" + tempfile.mkdtemp(prefix="lo_profile_")]
    if infilter:
        cmd = [lo, "--headless", f"--infilter={infilter}", "--convert-to", convert_to, "--outdir", output_dir]
        if fresh:
            cmd += ["-env:UserInstallation=file://" + tempfile.mkdtemp(prefix="lo_profile_")]
    cmd.append(input_path)

    subprocess.run(cmd, check=True, capture_output=True, text=True)

    default_out = os.path.join(output_dir, os.path.splitext(os.path.basename(input_path))[0] + "." + out_ext)
    if default_out != output_path:
        os.replace(default_out, output_path)


def convert_batch(input_paths: list[str], output_dir: str | None, fresh: bool = False):
    in_ext = os.path.splitext(input_paths[0])[1].lower()
    out_ext = TARGET_EXT.get(in_ext, "pdf")

    for p in input_paths:
        if os.path.splitext(p)[1].lower() != in_ext:
            print("Error: mixed file types not supported", file=sys.stderr)
            sys.exit(1)

    abs_paths = [os.path.abspath(p) for p in input_paths]

    out_dir = os.path.abspath(output_dir) if output_dir else get_default_dir()
    os.makedirs(out_dir, exist_ok=True)

    lo = find_libreoffice()

    convert_to = out_ext
    infilter = None
    if in_ext == ".pdf" and out_ext == "pptx":
        convert_to = "pptx:Impress Office Open XML"
        infilter = "impress_pdf_import"

    cmd = [lo, "--headless", "--convert-to", convert_to, "--outdir", out_dir]
    if fresh:
        cmd += ["-env:UserInstallation=file://" + tempfile.mkdtemp(prefix="lo_profile_")]
    if infilter:
        cmd = [lo, "--headless", f"--infilter={infilter}", "--convert-to", convert_to, "--outdir", out_dir]
        if fresh:
            cmd += ["-env:UserInstallation=file://" + tempfile.mkdtemp(prefix="lo_profile_")]
    cmd.extend(abs_paths)

    subprocess.run(cmd, check=True, capture_output=True, text=True)


def pdf_to_images(pdf_path: str, dpi: int = FAITHFUL_DPI) -> list[str]:
    out_dir = tempfile.mkdtemp(prefix="faithful_")
    prefix = os.path.join(out_dir, "slide")
    subprocess.run(
        ["pdftoppm", "-r", str(dpi), "-png", pdf_path, prefix],
        check=True, capture_output=True, text=True,
    )
    pngs = sorted(glob.glob(prefix + "-*.png"), key=lambda x: int(x.rsplit("-", 1)[1].split(".")[0]))
    return pngs


def ppdf_faithful(input_path: str, output_path: str):
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    output_dir = os.path.dirname(output_path)
    os.makedirs(output_dir, exist_ok=True)

    print("Rendering slides via LibreOffice...", file=sys.stderr)
    tmp_pdf = tempfile.mktemp(suffix=".pdf")
    pngs = []
    try:
        convert_single(input_path, tmp_pdf)
        print("Rasterizing pages...", file=sys.stderr)
        pngs = pdf_to_images(tmp_pdf, FAITHFUL_DPI)
        if not pngs:
            print("Error: no pages rendered", file=sys.stderr)
            sys.exit(1)
        print(f"Building PDF from {len(pngs)} slide image(s)...", file=sys.stderr)
        import img2pdf
        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(pngs, dpi=FAITHFUL_DPI))
    finally:
        if pngs:
            shutil.rmtree(os.path.dirname(pngs[0]), ignore_errors=True)
        if os.path.exists(tmp_pdf):
            os.unlink(tmp_pdf)


def pdf2pptx_faithful(input_path: str, output_path: str):
    input_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    output_dir = os.path.dirname(output_path)
    os.makedirs(output_dir, exist_ok=True)

    print("Rasterizing PDF pages...", file=sys.stderr)
    pngs = pdf_to_images(input_path, FAITHFUL_DPI)
    if not pngs:
        print("Error: no pages rendered", file=sys.stderr)
        sys.exit(1)

    print(f"Building PPTX with {len(pngs)} slide image(s)...", file=sys.stderr)
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from PIL import Image

    prs = Presentation()
    sw = prs.slide_width
    sh = prs.slide_height

    for i, png_path in enumerate(pngs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = Image.open(png_path)
        iw, ih = img.size
        scale = min(sw / iw, sh / ih)
        nw = int(iw * scale)
        nh = int(ih * scale)
        left = int((sw - nw) / 2)
        top = int((sh - nh) / 2)
        slide.shapes.add_picture(png_path, left, top, nw, nh)
        print(f"  Slide {i+1}/{len(pngs)}", file=sys.stderr)

    prs.save(output_path)


def convert_single_faithful(input_path: str, output_path: str):
    in_ext = os.path.splitext(input_path)[1].lower()
    out_ext = os.path.splitext(output_path)[1].lower().lstrip(".")
    if in_ext == ".pptx" and out_ext == "pdf":
        ppdf_faithful(input_path, output_path)
    elif in_ext == ".pdf" and out_ext == "pptx":
        pdf2pptx_faithful(input_path, output_path)
    else:
        print(f"Error: unsupported conversion: {in_ext} -> {out_ext}", file=sys.stderr)
        sys.exit(1)


def convert_batch_faithful(input_paths: list[str], output_dir: str | None):
    in_ext = os.path.splitext(input_paths[0])[1].lower()
    out_ext = TARGET_EXT.get(in_ext, "pdf")

    for p in input_paths:
        if os.path.splitext(p)[1].lower() != in_ext:
            print("Error: mixed file types not supported", file=sys.stderr)
            sys.exit(1)

    out_dir = os.path.abspath(output_dir) if output_dir else get_default_dir()
    os.makedirs(out_dir, exist_ok=True)

    for p in input_paths:
        out_name = os.path.splitext(os.path.basename(p))[0] + "." + out_ext
        out_path = os.path.join(out_dir, out_name)
        print(f"Converting: {p}", file=sys.stderr)
        print(f"Output:     {out_path}", file=sys.stderr)
        convert_single_faithful(p, out_path)


def check_fonts():
    common = ["Calibri", "Aptos", "Segoe UI", "Times New Roman", "Arial", "Verdana",
              "Trebuchet MS", "Tahoma", "Cambria", "Candara", "Consolas", "Constantia",
              "Corbel", "Franklin Gothic", "Gabriola"]
    try:
        fc_out = subprocess.run(["fc-list", ":lang=en"], capture_output=True, text=True, timeout=10).stdout.lower()
    except Exception:
        return
    missing = [f for f in common if f.lower() not in fc_out]
    if missing:
        print(f"Warning: common MS fonts not found: {', '.join(missing[:6])}...", file=sys.stderr)
        print(f"  Install with: sudo pacman -S ttf-liberation ttf-croscore", file=sys.stderr)
        print(f"  (Or install ttf-ms-fonts from AUR for exact MS fonts)", file=sys.stderr)


def main():
    import argparse

    parser = argparse.ArgumentParser(description="Convert between PPTX and PDF")
    parser.add_argument("inputs", nargs="+", help="Input file(s) (.pptx or .pdf)")
    parser.add_argument("output_dir", nargs="?", default=None, help="Output directory (default: ~/Documents)")
    parser.add_argument("-o", "--output", help="Output filename (only with single input)")
    parser.add_argument("--dpi", type=int, default=150, help="Rendering DPI (ignored, LibreOffice default used)")
    parser.add_argument("--delete", "-d", action="store_true", help="Delete original files after conversion")
    parser.add_argument("--fresh", action="store_true", help="Use fresh LibreOffice profile (fix footer/header issues)")
    parser.add_argument("--faithful", action="store_true",
                        help="Render slides as images for pixel-identical output (slower, larger files)")

    args = parser.parse_args()

    out_dir = args.output_dir
    remaining = list(args.inputs)
    if out_dir is None and len(remaining) > 1:
        last = remaining[-1]
        last_exp = os.path.expanduser(last)
        is_ext = os.path.splitext(last)[1].lower() in (".pptx", ".pdf")
        if not is_ext and (os.path.isdir(last_exp) or not glob.glob(last)):
            out_dir = last_exp
            remaining.pop()

    input_files = []
    for item in remaining:
        expanded = sorted(glob.glob(item))
        if not expanded:
            print(f"Error: no files match: {item}", file=sys.stderr)
            sys.exit(1)
        input_files.extend(expanded)

    for p in input_files:
        if not os.path.isfile(p):
            print(f"Error: not a file: {p}", file=sys.stderr)
            sys.exit(1)

        in_ext = os.path.splitext(p)[1].lower()
        if in_ext not in TARGET_EXT:
            print(f"Error: unsupported format: {in_ext} (use .pptx or .pdf)", file=sys.stderr)
            sys.exit(1)

    if args.output and len(input_files) > 1:
        print("Error: -o flag requires a single input file", file=sys.stderr)
        sys.exit(1)

    if args.faithful:
        check_fonts()
        try:
            if len(input_files) == 1:
                output = resolve_output(input_files[0], args.output, out_dir)
                print(f"Converting: {input_files[0]}", file=sys.stderr)
                print(f"Output:     {output}", file=sys.stderr)
                convert_single_faithful(input_files[0], output)
            else:
                print(f"Converting {len(input_files)} file(s)...", file=sys.stderr)
                convert_batch_faithful(input_files, out_dir)
        except (subprocess.CalledProcessError, FileNotFoundError, OSError) as e:
            print(f"Error: {e}", file=sys.stderr)
            sys.exit(1)
    else:
        try:
            if len(input_files) == 1:
                output = resolve_output(input_files[0], args.output, out_dir)
                print(f"Converting: {input_files[0]}", file=sys.stderr)
                print(f"Output:     {output}", file=sys.stderr)
                convert_single(input_files[0], output, fresh=args.fresh)
            else:
                print(f"Converting {len(input_files)} file(s)...", file=sys.stderr)
                convert_batch(input_files, out_dir, fresh=args.fresh)
        except (subprocess.CalledProcessError, FileNotFoundError, OSError) as e:
            print(f"Error: {e}", file=sys.stderr)
            sys.exit(1)

    if args.delete:
        for f in input_files:
            os.remove(f)
            print(f"Deleted: {f}", file=sys.stderr)


if __name__ == "__main__":
    main()
